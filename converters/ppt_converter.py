"""PowerPoint (.pptx/.ppt) 转 Markdown 转换器"""
import os
from pathlib import Path


def convert_ppt(input_path, output_dir, **kwargs):
    """将 PowerPoint 文件转换为 Markdown。

    Args:
        input_path: PPT文件路径
        output_dir: 输出目录

    Returns:
        转换后的Markdown文本
    """
    input_path = Path(input_path)
    ext = input_path.suffix.lower()

    if ext == '.pptx':
        return _convert_pptx(input_path, output_dir)
    elif ext == '.ppt':
        return _convert_ppt(input_path, output_dir)
    else:
        raise ValueError(f"不支持的 PPT 格式: {ext}")


def _convert_pptx(input_path, output_dir):
    """使用 python-pptx 转换 .pptx 文件。"""
    from pptx import Presentation
    from pptx.util import Inches
    import io

    prs = Presentation(str(input_path))
    md_parts = []
    images_dir = output_dir / f"{input_path.stem}_images"
    image_counter = 0

    md_parts.append(f"# {input_path.stem}\n")

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = []
        slide_parts.append(f"## 幻灯片 {slide_num}")

        # 提取幻灯片备注
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        for shape in slide.shapes:
            # 处理文本
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # 检测标题级别
                    if para.level == 0 and shape.shape_id <= 2:
                        slide_parts.append(f"### {text}")
                    elif para.level > 0:
                        indent = "  " * (para.level - 1)
                        slide_parts.append(f"{indent}- {text}")
                    else:
                        slide_parts.append(text)

            # 处理表格
            if shape.has_table:
                table = shape.table
                table_md = _table_to_markdown(table)
                slide_parts.append(table_md)

            # 处理图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_counter += 1
                    images_dir.mkdir(parents=True, exist_ok=True)
                    image = shape.image
                    ext = image.content_type.split('/')[-1]
                    if ext == 'jpeg':
                        ext = 'jpg'
                    img_filename = f"slide{slide_num}_img{image_counter}.{ext}"
                    img_path = images_dir / img_filename
                    img_path.write_bytes(image.blob)
                    slide_parts.append(
                        f"![幻灯片{slide_num}图片{image_counter}]"
                        f"({input_path.stem}_images/{img_filename})"
                    )
                except Exception:
                    pass

        # 添加备注
        if notes_text:
            slide_parts.append(f"\n> **备注:** {notes_text}")

        md_parts.append("\n\n".join(slide_parts))

    # 清理空的图片目录
    if images_dir.exists() and not any(images_dir.iterdir()):
        images_dir.rmdir()

    return "\n\n---\n\n".join(md_parts) + "\n"


def _table_to_markdown(table):
    """将 PPT 表格转换为 Markdown 表格。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip().replace('|', '\\|'))
        rows.append(cells)

    if not rows:
        return ""

    # 构建 Markdown 表格
    md_lines = []
    # 表头
    md_lines.append("| " + " | ".join(rows[0]) + " |")
    # 分隔线
    md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    # 数据行
    for row in rows[1:]:
        # 确保列数一致
        while len(row) < len(rows[0]):
            row.append("")
        md_lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")

    return "\n".join(md_lines)


def _convert_ppt(input_path, output_dir):
    """转换旧版 .ppt 文件（需要 LibreOffice）。"""
    import subprocess
    import tempfile

    libreoffice_paths = [
        'libreoffice',
        'soffice',
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
    ]

    lo_cmd = None
    for path in libreoffice_paths:
        try:
            subprocess.run([path, '--version'], capture_output=True, timeout=10)
            lo_cmd = path
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    if lo_cmd is None:
        raise RuntimeError(
            ".ppt 格式需要安装 LibreOffice 来转换。\n"
            "macOS 安装方法: brew install --cask libreoffice\n"
            "或者将 .ppt 文件手动另存为 .pptx 格式"
        )

    with tempfile.TemporaryDirectory() as temp_dir:
        subprocess.run(
            [lo_cmd, '--headless', '--convert-to', 'pptx', '--outdir', temp_dir, str(input_path)],
            capture_output=True,
            timeout=120,
        )

        pptx_path = Path(temp_dir) / f"{input_path.stem}.pptx"
        if pptx_path.exists():
            return _convert_pptx(pptx_path, output_dir)
        else:
            raise RuntimeError(f"LibreOffice 转换 {input_path.name} 失败")
